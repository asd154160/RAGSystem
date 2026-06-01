"""
统一文档解析器 - 支持 TXT/MD/PDF/DOCX/XLSX/PPTX
产出统一的 Block 中间结构
"""
import io
import os
import re
import uuid
import hashlib
from dataclasses import dataclass, field


@dataclass
class Block:
    block_id: str = field(default_factory=lambda: str(uuid.uuid4()))
    type: str = "paragraph"  # title, paragraph, table, list, image_caption
    text: str = ""
    page_no: int | None = None
    sheet_name: str | None = None
    slide_no: int | None = None
    section_title: str | None = None
    order_index: int = 0
    metadata: dict = field(default_factory=dict)


@dataclass
class ParseResult:
    document_id: str = ""
    document_version_id: str = ""
    blocks: list[Block] = field(default_factory=list)
    content_hash: str = ""


def _clean_text(text: str) -> str:
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' +', ' ', text)
    return text.strip()


def _parse_txt(content: str) -> list[Block]:
    blocks = []
    for para in content.split("\n\n"):
        para = _clean_text(para)
        if not para:
            continue
        bt = "title" if len(para) < 80 and not para.endswith("。") else "paragraph"
        blocks.append(Block(type=bt, text=para, order_index=len(blocks)))
    return blocks


def _parse_md(content: str) -> list[Block]:
    blocks = []
    lines = content.split("\n")
    current_section = None
    buf = []

    def flush():
        nonlocal buf
        text = _clean_text("\n".join(buf))
        buf.clear()
        if text:
            blocks.append(Block(type="paragraph", text=text, section_title=current_section, order_index=len(blocks)))

    for line in lines:
        stripped = line.strip()
        if line.startswith("# "):
            flush()
            current_section = stripped.lstrip("# ")
            blocks.append(Block(type="title", text=current_section, section_title=current_section, order_index=len(blocks)))
        elif line.startswith("## ") or line.startswith("### "):
            flush()
            current_section = stripped.lstrip("# ")
        elif stripped.startswith(("- ", "* ")):
            flush()
            blocks.append(Block(type="list", text=stripped, section_title=current_section, order_index=len(blocks)))
        elif stripped.startswith("|"):
            flush()
            blocks.append(Block(type="table", text=stripped, section_title=current_section, order_index=len(blocks)))
        elif stripped:
            buf.append(stripped)
        else:
            flush()
    flush()
    return blocks


def _parse_pdf(file_path: str) -> list[Block]:
    blocks = []
    try:
        import fitz
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc, 1):
            page_text = page.get_text()
            if page_text.strip():
                for para in page_text.split("\n\n"):
                    para = _clean_text(para)
                    if para:
                        blocks.append(Block(type="paragraph", text=para, page_no=page_num, order_index=len(blocks)))
            tables = page.find_tables()
            if tables:
                for t in tables:
                    ex = t.extract()
                    if ex and ex[0]:
                        rows = [" | ".join(str(c).strip() if c else "" for c in row) for row in ex]
                        blocks.append(Block(type="table", text="\n".join(rows), page_no=page_num, order_index=len(blocks)))
        doc.close()
        return blocks
    except ImportError:
        pass

    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text:
                for para in text.split("\n\n"):
                    para = _clean_text(para)
                    if para:
                        blocks.append(Block(type="paragraph", text=para, page_no=page_num, order_index=len(blocks)))
    except Exception:
        pass
    return blocks


def _parse_docx(file_path: str) -> list[Block]:
    from docx import Document
    blocks = []
    doc = Document(file_path)
    for para in doc.paragraphs:
        text = _clean_text(para.text)
        if not text:
            continue
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            blocks.append(Block(type="title", text=text, section_title=text, order_index=len(blocks)))
        else:
            blocks.append(Block(type="paragraph", text=text, order_index=len(blocks)))
    for table in doc.tables:
        rows_text = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        if rows_text:
            blocks.append(Block(type="table", text="\n".join(rows_text), order_index=len(blocks)))
    return blocks


def _parse_xlsx(file_path: str) -> list[Block]:
    from openpyxl import load_workbook
    blocks = []
    wb = load_workbook(file_path, data_only=True)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row_idx, row in enumerate(ws.iter_rows(values_only=True), 1):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if all(not c for c in cells):
                continue
            blocks.append(Block(type="table", text=" | ".join(cells), sheet_name=sheet_name, order_index=len(blocks)))
    wb.close()
    return blocks


def _parse_pptx(file_path: str) -> list[Block]:
    from pptx import Presentation
    blocks = []
    prs = Presentation(file_path)
    for slide_no, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        blocks.append(Block(type="paragraph" if slide_no > 1 and blocks else "title", text=t, slide_no=slide_no, order_index=len(blocks)))
            if shape.has_table:
                rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows]
                if rows:
                    blocks.append(Block(type="table", text="\n".join(rows), slide_no=slide_no, order_index=len(blocks)))
    return blocks


def parse_file(file_path: str) -> ParseResult:
    ext = file_path.rsplit(".", 1)[-1].lower()
    with open(file_path, "rb") as f:
        raw = f.read()
    content_hash = hashlib.sha256(raw).hexdigest()

    if ext == "txt":
        blocks = _parse_txt(raw.decode("utf-8", errors="replace"))
    elif ext == "md":
        blocks = _parse_md(raw.decode("utf-8", errors="replace"))
    elif ext == "pdf":
        blocks = _parse_pdf(file_path)
    elif ext == "docx":
        blocks = _parse_docx(file_path)
    elif ext == "xlsx":
        blocks = _parse_xlsx(file_path)
    elif ext == "pptx":
        blocks = _parse_pptx(file_path)
    else:
        content = raw.decode("utf-8", errors="replace")
        blocks = _parse_txt(content)

    return ParseResult(blocks=blocks, content_hash=content_hash)


def parse_from_bytes(file_name: str, data: bytes) -> ParseResult:
    ext = file_name.rsplit(".", 1)[-1].lower()
    content_hash = hashlib.sha256(data).hexdigest()

    if ext in ("pdf", "docx", "xlsx", "pptx"):
        tmp_path = f"/tmp/parse_{uuid.uuid4().hex}.{ext}"
        with open(tmp_path, "wb") as f:
            f.write(data)
        try:
            result = parse_file(tmp_path)
            result.content_hash = content_hash
            return result
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    text = data.decode("utf-8", errors="replace")
    blocks = _parse_txt(text) if ext == "txt" else _parse_md(text)
    return ParseResult(blocks=blocks, content_hash=content_hash)
